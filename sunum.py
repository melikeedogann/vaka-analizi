# ======================================================================
# ZİNCİR SÜPERMARKETLERİ POWERPOINT SUNUM OLUŞTURUCU
# ======================================================================
# Bu script, süpermarket satış verilerinin analizini PowerPoint sunumuna
# dönüştürür. Analiz sonuçlarını grafikler ve metinlerle görselleştirir.
# Professional bir sunum çıktısı oluşturur.
# ======================================================================

# KÜTÜPHANE İÇE AKTARIMLARI
# =========================
# Veri manipülasyonu ve analizi için pandas
import pandas as pd
# Grafik oluşturma için matplotlib ve seaborn
import matplotlib.pyplot as plt
import seaborn as sns
# Matematiksel işlemler için numpy
import numpy as np
# PowerPoint sunum oluşturma için python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# VERİ YÜKLEME VE ÖN İŞLEME
# =========================
# CSV dosyasından veri setini yükle
df = pd.read_csv('2425 vaka02 - zincir supermarket_sales - Sheet1.csv')

# TARİH VERİLERİNİN İŞLENMESİ
# ===========================
# Tarih sütununu datetime formatına dönüştür
df['Date'] = pd.to_datetime(df['Date'])
# Tarihten ay, gün ve haftanın günü bilgilerini çıkar
df['Month'] = df['Date'].dt.month
df['Day'] = df['Date'].dt.day
df['DayOfWeek'] = df['Date'].dt.dayofweek

# TEMEL ANALİZ SONUÇLARINI HESAPLAMA
# ==================================
# Bu bölüm sunumda kullanılacak tüm analiz sonuçlarını önceden hesaplar

# ÜRÜN KATEGORİLERİ ANALİZİ
# Ürün kategorilerine göre detaylı istatistiksel analiz
product_revenue = df.groupby('Product line').agg({
    'Total': 'sum',           # Toplam satış tutarı
    'gross income': 'sum',    # Toplam brüt gelir
    'Invoice ID': 'count',    # İşlem sayısı
    'Quantity': 'sum',        # Toplam satılan ürün adedi
    'Unit price': 'mean',     # Ortalama birim fiyat
    'Rating': 'mean'          # Ortalama müşteri memnuniyeti
}).sort_values('Total', ascending=False)

# Her kategorinin toplam gelir içindeki yüzde payını hesapla
product_revenue['Revenue_Percentage'] = product_revenue['Total'] / product_revenue['Total'].sum() * 100

# MÜŞTERİ TİPLERİ ANALİZİ
# Müşteri tiplerine göre detaylı istatistiksel analiz
customer_revenue = df.groupby('Customer type').agg({
    'Total': 'sum',
    'gross income': 'sum',
    'Invoice ID': 'count',
    'Quantity': 'sum'
}).sort_values('Total', ascending=False)

customer_revenue['Revenue_Percentage'] = customer_revenue['Total'] / customer_revenue['Total'].sum() * 100

# ŞUBE PERFORMANS ANALİZİ
# Şubelere göre performans karşılaştırması
branch_revenue = df.groupby('Branch').agg({
    'Total': 'sum',
    'gross income': 'sum',
    'Invoice ID': 'count',
    'Quantity': 'sum'
}).sort_values('Total', ascending=False)

branch_revenue['Revenue_Percentage'] = branch_revenue['Total'] / branch_revenue['Total'].sum() * 100

# CİNSİYET BAZINDA ANALİZ
# Erkek ve kadın müşterilerin harcama davranışları
gender_revenue = df.groupby('Gender').agg({
    'Total': 'sum',
    'gross income': 'sum',
    'Invoice ID': 'count',
    'Quantity': 'sum'
}).sort_values('Total', ascending=False)

gender_revenue['Revenue_Percentage'] = gender_revenue['Total'] / gender_revenue['Total'].sum() * 100

# ÖDEME YÖNTEMİ ANALİZİ
# Ödeme yöntemlerine göre gelir analizi
payment_revenue = df.groupby('Payment').agg({
    'Total': 'sum',
    'gross income': 'sum',
    'Invoice ID': 'count',
    'Quantity': 'sum'
}).sort_values('Total', ascending=False)

payment_revenue['Revenue_Percentage'] = payment_revenue['Total'] / payment_revenue['Total'].sum() * 100

# MÜŞTERİ MEMNUNİYETİ ANALİZİ
# Ürün kategorilerine göre müşteri memnuniyet seviyesi
rating_by_product = df.groupby('Product line')['Rating'].mean().sort_values(ascending=False)

# ======================================================================
# POWERPOINT SUNUM OLUŞTURMA BÖLÜMÜ
# ======================================================================
# Yeni bir PowerPoint sunumu oluştur
prs = Presentation()

# SLAYT 1: KAPAK SAYFASI
# ======================
# Boş bir kapak slaytı oluştur
slide = prs.slides.add_slide(prs.slide_layouts[0])  # Başlık slaytı düzeni
title = slide.shapes.title                          # Başlık alanına erişim
subtitle = slide.placeholders[1]                    # Alt başlık alanına erişim

# Kapak sayfası metinlerini ayarla
title.text = "ZİNCİR SÜPERMARKETLERİ VERİ ANALİZİ"
subtitle.text = "3 Şubede 3 Aylık Satış Verilerinin Analizi"

# SLAYT 2: İÇERİK SAYFASI
# =======================
# İçerik tablosu slaytı oluştur
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Başlık ve içerik düzeni
title = slide.shapes.title
title.text = "İçerik"

# İçerik metnini hazırla
content = slide.placeholders[1]
content.text = """
1. Analiz Yöntemi ve Veri Seti
2. İşletme Gelirlerinin Analizi
   - Ürün Kategorilerine Göre
   - Müşteri Tiplerine Göre
   - Şube Bazında
3. Müşteri Segmentasyonu
4. Müşteri Memnuniyeti
5. Önerilen Stratejiler
"""

# SLAYT 3: VERİ SETİ BİLGİLERİ
# ============================
# Veri seti ve analiz yöntemi açıklama slaytı
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Veri Seti ve Analiz Yöntemi"

# Veri seti özellikleri ve analiz yöntemi bilgileri
content = slide.placeholders[1]
content.text = """
Veri Seti Özellikleri:
• 1000 adet satış kaydı
• 3 farklı şubede 3 aylık satış verileri (Ocak-Mart 2019)
• 17 değişken: Fatura ID, Şube, Şehir, Müşteri Tipi, Cinsiyet, Ürün Kategorisi, vb.

Analiz Yöntemi:
• Tanımlayıcı İstatistikler
• Karşılaştırmalı Analizler
• Gelir ve Karlılık Analizleri
• Müşteri Segmentasyonu
• Trend Analizleri
"""

# SLAYT 4: ÜRÜN KATEGORİLERİ GELİR ANALİZİ
# =========================================
# Sadece grafik gösteren slayt düzeni
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Ürün Kategorilerine Göre Gelir Dağılımı"

# ÜRÜN KATEGORİLERİ GRAFİĞİ OLUŞTURMA
# Matplotlib ile ürün kategorileri gelir grafiği oluştur
plt.figure(figsize=(10, 6))
product_revenue['Total'].plot(kind='bar', color='skyblue')
plt.title('Ürün Kategorilerine Göre Toplam Gelir')
plt.xlabel('Ürün Kategorisi')
plt.ylabel('Toplam Gelir ($)')
plt.xticks(rotation=45, ha='right')  # X ekseni etiketlerini döndür

# Her çubuğun üstüne yüzde bilgisini ekle
for i, v in enumerate(product_revenue['Revenue_Percentage']):
    plt.text(i, product_revenue['Total'].iloc[i] + 1000, f'%{v:.1f}', ha='center')

plt.tight_layout()
# Grafiği PNG dosyası olarak kaydet
plt.savefig('urun_gelir.png', dpi=300, bbox_inches='tight')
plt.close()  # Bellek tasarrufu için grafiği kapat

# Grafiği PowerPoint slaytına ekle
left = Inches(0.5)    # Soldan mesafe
top = Inches(1.5)     # Üstten mesafe
width = Inches(9)     # Grafik genişliği
height = Inches(5)    # Grafik yüksekliği
slide.shapes.add_picture('urun_gelir.png', left, top, width, height)

# SLAYT 5: MÜŞTERİ TİPLERİ GELİR ANALİZİ
# =======================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Müşteri Tiplerine Göre Gelir Dağılımı"

# MÜŞTERİ TİPİ GRAFİĞİ OLUŞTURMA
plt.figure(figsize=(10, 6))
customer_revenue['Total'].plot(kind='bar', color='lightgreen')
plt.title('Müşteri Tiplerine Göre Toplam Gelir')
plt.xlabel('Müşteri Tipi')
plt.ylabel('Toplam Gelir ($)')

# Yüzde bilgilerini çubukların üstüne ekle
for i, v in enumerate(customer_revenue['Revenue_Percentage']):
    plt.text(i, customer_revenue['Total'].iloc[i] + 1000, f'%{v:.1f}', ha='center')

plt.tight_layout()
plt.savefig('musteri_gelir.png', dpi=300, bbox_inches='tight')
plt.close()

# Grafiği slayta yerleştir
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(5)
slide.shapes.add_picture('musteri_gelir.png', left, top, width, height)

# SLAYT 6: ŞUBE PERFORMANS ANALİZİ
# ================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Şubelere Göre Gelir Dağılımı"

# ŞUBE PERFORMANS GRAFİĞİ OLUŞTURMA
plt.figure(figsize=(10, 6))
branch_revenue['Total'].plot(kind='bar', color='salmon')
plt.title('Şubelere Göre Toplam Gelir')
plt.xlabel('Şube')
plt.ylabel('Toplam Gelir ($)')

# Yüzde değerlerini ekle
for i, v in enumerate(branch_revenue['Revenue_Percentage']):
    plt.text(i, branch_revenue['Total'].iloc[i] + 1000, f'%{v:.1f}', ha='center')

plt.tight_layout()
plt.savefig('sube_gelir.png', dpi=300, bbox_inches='tight')
plt.close()

# Grafiği slayta yerleştir
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(5)
slide.shapes.add_picture('sube_gelir.png', left, top, width, height)

# SLAYT 7: ŞUBE-ÜRÜN KATEGORİSİ MATRİSİ
# ======================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Şube-Ürün Kategorisi Matrisi"

# ŞUBE-ÜRÜN MATRİSİ GRAFİĞİ OLUŞTURMA
# Şube ve ürün kategorisi kombinasyonu için pivot tablo oluştur
branch_product_pivot = pd.pivot_table(df, values='Total',
                                    index='Branch',
                                    columns='Product line',
                                    aggfunc='sum')

plt.figure(figsize=(12, 6))
# Isı haritası (heatmap) oluştur - hangi şubede hangi ürün kategorisinin daha başarılı olduğunu gösterir
sns.heatmap(branch_product_pivot, annot=True, fmt='.0f', cmap='YlGnBu', linewidths=0.5)
plt.title('Şube-Ürün Kategorisi Gelir Matrisi')
plt.ylabel('Şube')
plt.xlabel('Ürün Kategorisi')
plt.tight_layout()
plt.savefig('sube_urun_matrisi.png', dpi=300, bbox_inches='tight')
plt.close()

# Grafiği slayta yerleştir
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(5)
slide.shapes.add_picture('sube_urun_matrisi.png', left, top, width, height)

# SLAYT 8: MÜŞTERİ MEMNUNİYETİ ANALİZİ
# =====================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Ürün Kategorilerine Göre Müşteri Memnuniyeti"

# MÜŞTERİ MEMNUNİYETİ GRAFİĞİ OLUŞTURMA
plt.figure(figsize=(10, 6))
rating_by_product.plot(kind='bar', color='gold')
plt.title('Ürün Kategorilerine Göre Ortalama Değerlendirme Puanı')
plt.xlabel('Ürün Kategorisi')
plt.ylabel('Ortalama Puan (1-10)')
plt.xticks(rotation=45, ha='right')
plt.ylim(0, 10)  # Y ekseni limitini 0-10 arasında ayarla
plt.tight_layout()
plt.savefig('urun_memnuniyet.png', dpi=300, bbox_inches='tight')
plt.close()

# Grafiği slayta yerleştir
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(5)
slide.shapes.add_picture('urun_memnuniyet.png', left, top, width, height)

# SLAYT 9: MÜŞTERİ SEGMENTASYONU ANALİZİ
# ======================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Müşteri Segmentasyonu Analizi"

# MÜŞTERİ SEGMENTASYON GRAFİĞİ OLUŞTURMA
# Müşteri tipi ve cinsiyet kombinasyonu analizi
customer_gender_revenue = df.groupby(['Customer type', 'Gender'])['Total'].sum().unstack()
plt.figure(figsize=(10, 6))
customer_gender_revenue.plot(kind='bar', color=['skyblue', 'salmon'])
plt.title('Müşteri Tipi ve Cinsiyet Bazında Toplam Gelir')
plt.xlabel('Müşteri Tipi')
plt.ylabel('Toplam Gelir ($)')
plt.legend(title='Cinsiyet')  # Renk kodları için açıklama
plt.tight_layout()
plt.savefig('musteri_cinsiyet_gelir.png', dpi=300, bbox_inches='tight')
plt.close()

# Grafiği slayta yerleştir
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(5)
slide.shapes.add_picture('musteri_cinsiyet_gelir.png', left, top, width, height)

# SLAYT 10: ZAMAN BAZLI ANALİZ
# ============================
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Zaman Bazlı Analiz"

# SAATLİK ANALİZ GRAFİĞİ OLUŞTURMA
# Saatlik satış trendlerini analiz et
hourly_analysis = df.copy()
hourly_analysis['Hour'] = pd.to_datetime(df['Time']).dt.hour
hourly_revenue = hourly_analysis.groupby('Hour').agg({
    'Total': 'sum'
}).sort_values('Total', ascending=False)

plt.figure(figsize=(10, 6))
# Çizgi grafik - zaman serileri için en uygun görselleştirme
hourly_revenue['Total'].sort_index().plot(kind='line', marker='o', color='teal')
plt.title('Saat Bazında Gelir Dağılımı')
plt.xlabel('Saat')
plt.ylabel('Toplam Gelir ($)')
plt.grid(True, linestyle='--', alpha=0.7)  # Grid çizgileri ekle
plt.tight_layout()
plt.savefig('saat_gelir.png', dpi=300, bbox_inches='tight')
plt.close()

# Grafiği slayta yerleştir
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(5)
slide.shapes.add_picture('saat_gelir.png', left, top, width, height)

# SLAYT 11: ÖDEME YÖNTEMLERİ ANALİZİ
# ===================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Ödeme Yöntemlerine Göre Gelir Dağılımı"

# ÖDEME YÖNTEMLERİ GRAFİĞİ OLUŞTURMA
plt.figure(figsize=(10, 6))
payment_revenue['Total'].plot(kind='bar', color='lightpink')
plt.title('Ödeme Yöntemlerine Göre Toplam Gelir')
plt.xlabel('Ödeme Yöntemi')
plt.ylabel('Toplam Gelir ($)')

# Yüzde bilgilerini çubukların üstüne ekle
for i, v in enumerate(payment_revenue['Revenue_Percentage']):
    plt.text(i, payment_revenue['Total'].iloc[i] + 1000, f'%{v:.1f}', ha='center')

plt.tight_layout()
plt.savefig('odeme_gelir.png', dpi=300, bbox_inches='tight')
plt.close()

# Grafiği slayta yerleştir
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(5)
slide.shapes.add_picture('odeme_gelir.png', left, top, width, height)

# SLAYT 12: ÖNERİLEN STRATEJİLER
# ==============================
# Metin tabanlı strateji önerileri slaytı
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Başlık ve içerik düzeni
title = slide.shapes.title
title.text = "Önerilen Stratejiler"

# Kapsamlı strateji önerileri metni
content = slide.placeholders[1]
content.text = """
1. Ürün Stratejileri:
   • Yiyecek ve İçecekler kategorisini güçlendirme
   • Spor ve Seyahat kategorisinde ürün çeşitliliği artırma
   • Sağlık ve Güzellik kategorisinin geliştirilmesi

2. Müşteri Stratejileri:
   • Üyelik programını güçlendirme
   • Erkek müşterilere yönelik kampanyalar
   • Müşteri sadakat programları

3. Şube Stratejileri:
   • Şube C'nin başarı faktörlerini diğer şubelere yayma
   • Şube B'de ürün gamı optimizasyonu
   • Şubeler arası performans karşılaştırma sistemi

4. Operasyonel Stratejiler:
   • Yoğun saatlerde personel sayısını artırma
   • E-cüzdan kullanımını teşvik etme
   • Stok yönetimi optimizasyonu
"""

# SLAYT 13: TEŞEKKÜR SAYFASI
# ==========================
# Kapanış slaytı oluştur
slide = prs.slides.add_slide(prs.slide_layouts[0])  # Başlık slaytı düzeni
title = slide.shapes.title
subtitle = slide.placeholders[1]

# Teşekkür mesajları
title.text = "TEŞEKKÜRLER"
subtitle.text = "Sorularınız?"

# ======================================================================
# SUNUM DOSYASINI KAYDETME
# ======================================================================
# PowerPoint sunumunu dosya olarak kaydet
prs.save('ZincirSupermarketleri_Analiz_Sunumu.pptx') 